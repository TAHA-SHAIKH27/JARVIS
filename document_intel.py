"""
document_intel.py
------------------
Plain-text extraction for uploaded documents, used by the document Q&A /
summarization pipeline (main.py's /api/document/extract + agent.py's
stream_document_analysis()). Supports PDF, PPTX, DOCX, TXT, and MD.
"""

import os

# Keep well under Gemini's context window while avoiding pathologically
# large payloads from someone uploading a 2000-page PDF.
MAX_CHARS = 400_000
MAX_PDF_PAGES = 200

SUPPORTED_EXTENSIONS = (".txt", ".md", ".csv", ".log", ".pdf", ".pptx", ".docx")


def extract_text(filepath: str, filename: str) -> dict:
    """Extract plain text from a document file on disk.

    Returns {"status": "success", "text": str, "char_count": int, "truncated": bool}
    or {"status": "error", "message": str}.
    """
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext in (".txt", ".md", ".csv", ".log"):
            text = _extract_txt(filepath)
        elif ext == ".pdf":
            text = _extract_pdf(filepath)
        elif ext == ".pptx":
            text = _extract_pptx(filepath)
        elif ext == ".docx":
            text = _extract_docx(filepath)
        else:
            return {
                "status": "error",
                "message": f"Unsupported file type '{ext or 'unknown'}', sir. I can read PDF, DOCX, PPTX, TXT, and MD files."
            }
    except Exception as e:
        return {"status": "error", "message": f"Failed to read {filename}, sir: {str(e)}"}

    text = (text or "").strip()
    if not text:
        return {
            "status": "error",
            "message": f"I couldn't find any readable text in {filename}, sir. It may be a scanned/image-only document."
        }

    truncated = False
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]
        truncated = True

    return {"status": "success", "text": text, "char_count": len(text), "truncated": truncated}


def _extract_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_pdf(filepath: str) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            if i >= MAX_PDF_PAGES:
                pages.append(f"\n[...truncated after {MAX_PDF_PAGES} pages...]")
                break
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(f"--- Page {i + 1} ---\n{page_text}")
    return "\n\n".join(pages)


def _extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        lines = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        lines.append(line)
            # Tables often hold key slide content too
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                    if row_text.strip():
                        lines.append(row_text)
        if lines:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _extract_docx(filepath: str) -> str:
    import docx
    doc = docx.Document(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)
